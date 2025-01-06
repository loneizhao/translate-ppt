import boto3
import json
from pptx import Presentation
import time
import threading
from datetime import datetime
import logging
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import os
from flask import Flask, request, render_template, send_file
from werkzeug.utils import secure_filename



logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)
translation_progress = {'current': 0, 'total': 0}

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
app.config['MAX_CONTENT_LENGTH'] = 160 * 1024 * 1024  # 16MB max file size

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# Create upload folder if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

class TokenBucket:
    def __init__(self, requests_per_minute):
        self.capacity = requests_per_minute
        self.tokens = requests_per_minute
        self.last_update = time.time()
        self.lock = threading.Lock()
        self.rate = requests_per_minute / 60.0  # tokens per second

    def get_token(self):
        with self.lock:
            now = time.time()
            time_passed = now - self.last_update
            self.tokens = min(self.capacity, self.tokens + time_passed * self.rate)
            self.last_update = now

            if self.tokens < 1:
                sleep_time = (1 - self.tokens) / self.rate
                time.sleep(sleep_time)
                self.tokens = 0
                self.last_update = time.time()
            else:
                self.tokens -= 1

class ThrottledPPTTranslator:
    def __init__(self, requests_per_minute=10, max_retries=3):
        self.bedrock = boto3.client(
            service_name='bedrock-runtime',
            region_name='us-east-1'
        )
        self.throttler = TokenBucket(requests_per_minute=requests_per_minute)
        self.translation_cache = {}
        self.error_log = []
        self.model_id = 'us.anthropic.claude-3-5-haiku-20241022-v1:0'
        self.max_retries = max_retries

    def set_consistent_font(self, run, target_language):
        """设置合适的字体以支持中英文"""
        if target_language.lower() in ['chinese', '中文']:
            compatible_fonts = [
                'Microsoft YaHei',  # 微软雅黑
                'SimSun',          # 宋体
                'SimHei',          # 黑体
                'DengXian',        # 等线体
                'Arial Unicode MS',
                'Source Han Sans', # 思源黑体
            ]
            
            current_font = run.font.name
            if current_font not in compatible_fonts:
                run.font.name = 'Microsoft YaHei'

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=4, max=10),
        retry=retry_if_exception_type((boto3.exceptions.Boto3Error, Exception)),
        after=lambda retry_state: logger.info(f"Retry attempt {retry_state.attempt_number}")
    )
    def translate_texts(self, texts, target_language, batch=True):
        """
        Unified translation method that handles both single and batch translations
        """
        if not texts:
            return [] if batch else ""

        try:
            self.throttler.get_token()
            
            if batch:
                combined_text = "\n---\n".join(texts)
            else:
                combined_text = texts

            system = [{
                "text": ("You are a professional translator. Please follow these rules:\n"
                        "1. Translate the text accurately while maintaining the original meaning\n"
                        "2. Preserve all original formatting, including line breaks and spacing\n"
                        "3. Keep technical terms, product names, and abbreviations in English (e.g., AWS, Amazon, EC2, API, SDK)\n"
                        "4. Do not add explanations or additional content\n"
                        "5. Maintain any special characters or symbols as they appear in the source text\n"
                        "6. Keep numbers, dates, and units in their original format")
            }]

            messages = [
                {"role": "user", "content": [{"text": f"Translate the following text to {target_language}. Provide only the translation without any additional comments or explanations:\n\n{combined_text}"}]}
            ]
            
            body = {
                "messages": messages,
                "system": system,
                "inferenceConfig": {
                    "maxTokens": 4096,
                    "topP": 0.1,
                    "temperature": 0.01
                }
            }

            response = self.bedrock.converse(
                modelId=self.model_id,
                **body
            )

            if 'output' in response and 'message' in response['output']:
                message_content = response['output']['message']['content']
                if message_content and len(message_content) > 0:
                    response_text = message_content[0]['text'].strip()
                    
                    if batch:
                        translations = [t.strip() for t in response_text.split('---')]
                        logger.info(f"Batch translations: {translations}")
                        
                        if len(translations) != len(texts):
                            logger.warning(f"Mismatch in translation count. Expected {len(texts)}, got {len(translations)}")
                            while len(translations) < len(texts):
                                translations.append(texts[len(translations)])
                            translations = translations[:len(texts)]
                        
                        return translations
                    else:
                        return response_text

            return texts if batch else texts

        except Exception as e:
            error_msg = "Batch translation error: " if batch else "Translation error: "
            self.error_log.append(f"{error_msg}{str(e)}")
            logger.error(f"{error_msg}{str(e)}")
            raise

    def translate_presentation_with_batching(self, input_file, target_language, batch_size=5):
        global translation_progress
        try:
            filename, ext = os.path.splitext(input_file)
            output_file = f"{filename}_cn{ext}"

            prs = Presentation(input_file)
            text_batch = []
            text_locations = []
            total_processed = 0
            failed_batches = []
            
            # Count total items
            total_items = sum(1 for slide in prs.slides 
                            for shape in slide.shapes 
                            if hasattr(shape, "text_frame")
                            for paragraph in shape.text_frame.paragraphs
                            for run in paragraph.runs
                            if run.text.strip())

            # Set initial progress
            translation_progress['total'] = total_items
            translation_progress['current'] = 0

            logger.info(f"Starting translation of {total_items} text elements")

            for slide_index, slide in enumerate(prs.slides):
                logger.info(f"Processing slide {slide_index + 1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    text_batch.append(run.text)
                                    text_locations.append((slide, shape, paragraph, run))
                                    
                                    if len(text_batch) >= batch_size:
                                        try:
                                            translations = self.translate_texts(text_batch, target_language, batch=True)
                                            for i, translation in enumerate(translations):
                                                if i < len(text_locations):
                                                    run = text_locations[i][3]
                                                    run.text = translation
                                                    self.set_consistent_font(run, target_language)
                                                    total_processed += 1
                                                    translation_progress['current'] = total_processed
                                                    logger.info(f"Progress: {total_processed}/{total_items}")
                                            text_batch = []
                                            text_locations = []
                                        except Exception as e:
                                            logger.error(f"Failed to translate batch: {e}")
                                            failed_batches.append((text_batch.copy(), text_locations.copy()))
                                            text_batch = []
                                            text_locations = []

            # Process remaining text
            if text_batch:
                try:
                    translations = self.translate_texts(text_batch, target_language, batch=True)
                    for i, translation in enumerate(translations):
                        if i < len(text_locations):
                            run = text_locations[i][3]
                            run.text = translation
                            self.set_consistent_font(run, target_language)
                            total_processed += 1
                except Exception as e:
                    logger.error(f"Failed to translate final batch: {e}")
                    failed_batches.append((text_batch.copy(), text_locations.copy()))

            prs.save(output_file)
            return output_file

        except Exception as e:
            logger.error(f"Error translating presentation: {str(e)}")
            raise
# Initialize translator
translator = ThrottledPPTTranslator()

@app.route('/progress')
def get_progress():
    return json.dumps({
        'current': translation_progress['current'],
        'total': translation_progress['total']
    })

@app.route('/')
def index():
    return '''
    <!doctype html>
    <html>
    <head>
        <title>PPT Translator</title>
        <!-- 添加 Font Awesome 图标库 -->
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/5.15.4/css/all.min.css">
        <!-- 添加 Bootstrap CSS -->
        <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
        <style>
            body {
                padding: 20px;
                background-color: #f8f9fa;
            }
            .container {
                max-width: 800px;
                margin: 0 auto;
                background-color: white;
                padding: 30px;
                border-radius: 10px;
                box-shadow: 0 0 10px rgba(0,0,0,0.1);
            }
            .upload-form {
                margin-top: 20px;
            }
            .file-input-wrapper {
                margin-bottom: 20px;
            }
            .progress {
                display: none;
                margin-top: 20px;
            }
            .result-section {
                margin-top: 20px;
                display: none;
            }
            .btn-translate {
                margin-top: 10px;
            }
            .download-btn {
                text-decoration: none;
                color: white;
            }
            .download-btn:hover {
                color: white;
            }
            .progress {
                height: 25px;
                margin-top: 20px;
                margin-bottom: 20px;
            }
            .progress-bar {
                transition: width 0.5s ease-in-out;
                text-align: center;
                line-height: 25px;
                color: white;
                font-weight: bold;
            }
        </style>
    </head>
    <body>
        <div class="container">
            <h1 class="text-center mb-4">
                <i class="fas fa-language"></i> PPT Translator
            </h1>
            
            <form action="/translate" method="post" enctype="multipart/form-data" class="upload-form" id="uploadForm">
                <div class="file-input-wrapper">
                    <label for="file" class="form-label">Select PPT File:</label>
                    <input type="file" class="form-control" name="file" id="file" accept=".ppt,.pptx" required>
                </div>
                <div class="progress">
                <div class="progress-bar progress-bar-striped progress-bar-animated" 
                        role="progressbar" 
                        aria-valuenow="0" 
                        aria-valuemin="0" 
                        aria-valuemax="100" 
                        style="width: 0%">
                        0%
                    </div>
                </div>
                <div class="mb-3">
                    <label for="target_language" class="form-label">Target Language:</label>
                    <select class="form-select" name="target_language" id="target_language">
                        <option value="chinese">Chinese</option>
                        <option value="japanese">English</option>
                    </select>
                </div>
                
                <div class="mb-3">
                    <label for="model" class="form-label">Select Model:</label>
                    <select class="form-select" name="model" id="model">
                        <option value="us.anthropic.claude-3-5-sonnet-20241022-v2:0">Claude 3 Sonnet</option>
                        <option value="us.anthropic.claude-3-5-haiku-20241022-v1:0">Claude 3 Haiku</option>
                        <option value="us.amazon.nova-pro-v1:0">Nova Pro</option>
                        
                    </select>
                </div>

                <button type="submit" class="btn btn-primary btn-translate w-100">
                    <i class="fas fa-translate"></i> Translate
                </button>
            </form>
            
            <div class="progress">
                <div class="progress-bar progress-bar-striped progress-bar-animated" role="progressbar" style="width: 0%"></div>
            </div>
            
            <div class="result-section" id="resultSection">
                <div class="alert alert-success">
                    <i class="fas fa-check-circle"></i> Translation completed!
                </div>
                <a href="#" class="btn btn-success w-100 download-btn" id="downloadBtn">
                    <i class="fas fa-download"></i> Download Translated File
                </a>
            </div>
        </div>

        <!-- 添加 Bootstrap JS 和其依赖 -->
        <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/js/bootstrap.bundle.min.js"></script>
        <script>
            document.getElementById('uploadForm').onsubmit = function(e) {
                e.preventDefault();
                
                const formData = new FormData(this);
                const progressBar = document.querySelector('.progress-bar');
                const progress = document.querySelector('.progress');
                const resultSection = document.getElementById('resultSection');
                const downloadBtn = document.getElementById('downloadBtn');
                
                progress.style.display = 'block';
                resultSection.style.display = 'none';
                
                let progressInterval = setInterval(function() {
                    fetch('/progress')
                        .then(response => response.json())
                        .then(data => {
                            if (data.total > 0) {
                                const percentage = (data.current / data.total) * 100;
                                progressBar.style.width = percentage + '%';
                                progressBar.setAttribute('aria-valuenow', percentage);
                                progressBar.textContent = Math.round(percentage) + '%';
                            }
                        });
                }, 1000);

                fetch('/translate', {
                    method: 'POST',
                    body: formData
                })
                .then(response => {
                    clearInterval(progressInterval);
                    
                    if (!response.ok) {
                        return response.json().then(data => {
                            throw new Error(data.error || 'Translation failed');
                        });
                    }
                    
                    // Get the filename from the Content-Disposition header
                    const contentDisposition = response.headers.get('Content-Disposition');
                    const filenameMatch = contentDisposition && contentDisposition.match(/filename="(.+)"/);
                    const filename = filenameMatch ? filenameMatch[1] : 'translated_file.pptx';
                    
                    return response.blob().then(blob => ({blob, filename}));
                })
                .then(({blob, filename}) => {
                    const url = window.URL.createObjectURL(blob);
                    downloadBtn.href = url;
                    downloadBtn.download = filename;
                    
                    progress.style.display = 'none';
                    resultSection.style.display = 'block';
                    
                    // Automatically trigger download
                    const tempLink = document.createElement('a');
                    tempLink.href = url;
                    tempLink.download = filename;
                    document.body.appendChild(tempLink);
                    tempLink.click();
                    document.body.removeChild(tempLink);
                    
                    // Clean up the blob URL after a delay
                    setTimeout(() => {
                        window.URL.revokeObjectURL(url);
                    }, 1000);
                })
                .catch(error => {
                    clearInterval(progressInterval);
                    alert('Error: ' + error.message);
                    progress.style.display = 'none';
                });
            };

        </script>
    </body>
    </html>
    '''
@app.route('/translate', methods=['POST'])
def translate_file():
    input_path = None
    output_file = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        target_language = request.form.get('target_language', 'chinese')
        selected_model = request.form.get('model', 'us.anthropic.claude-3-haiku-20240307-v1:0')

        translator.model_id = selected_model
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
            
        if not file.filename.endswith(('.ppt', '.pptx')):
            return jsonify({'error': 'Invalid file type'}), 400
            
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = secure_filename(f"{timestamp}_{file.filename}")
        input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        file.save(input_path)
        
        output_file = translator.translate_presentation_with_batching(input_path, target_language)
        
        if not os.path.exists(output_file):
            raise FileNotFoundError("Translation output file not found")

        # Clean up input file only
        if input_path and os.path.exists(input_path):
            os.remove(input_path)
            logger.info(f"Cleaned up input file: {input_path}")

        # Return the translated file
        return send_file(
            output_file,
            as_attachment=True,
            download_name=f"translated_{file.filename}",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        logger.error(f"Translation error: {str(e)}", exc_info=True)
        # Clean up files in case of error
        try:
            if input_path and os.path.exists(input_path):
                os.remove(input_path)
            if output_file and os.path.exists(output_file):
                os.remove(output_file)
        except Exception as cleanup_error:
            logger.error(f"Error during cleanup: {cleanup_error}")
        return jsonify({'error': str(e)}), 500


    
    finally:
        # Clean up temporary files
        try:
            if input_path and os.path.exists(input_path):
                os.remove(input_path)
                logger.info(f"Cleaned up input file: {input_path}")
            if output_file and os.path.exists(output_file):
                os.remove(output_file)
                logger.info(f"Cleaned up output file: {output_file}")
        except Exception as e:
            logger.error(f"Error cleaning up files: {e}", exc_info=True)

            
if __name__ == '__main__':
    app.run(debug=True)
